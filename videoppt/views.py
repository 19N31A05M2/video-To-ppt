from ast import Not
import django

from django.http import HttpResponse
from django.shortcuts import render
from django.http import HttpResponse
from django.views.decorators.csrf import csrf_exempt
import collections.abc
import cv2 
import cv2 as cv
import os
from pptx import Presentation

old_frame = None
currentframe = 0
vidcap = None
# Create your views here.
def home(request):
    return render(request,"home.html")

def capture(base,test):

    hsv_base = cv.cvtColor(base, cv.COLOR_BGR2HSV)
    hsv_test = cv.cvtColor(test, cv.COLOR_BGR2HSV)


    h_bins = 50
    s_bins = 60
    histSize = [h_bins, s_bins]
    h_ranges = [0, 180]
    s_ranges = [0, 256]
    ranges = h_ranges + s_ranges
    channels = [0, 1]

    hist_base = cv.calcHist([hsv_base], channels, None, histSize, ranges, accumulate=False)
    cv.normalize(hist_base, hist_base, alpha=0, beta=1, norm_type=cv.NORM_MINMAX)
    hist_test = cv.calcHist([hsv_test], channels, None, histSize, ranges, accumulate=False)
    cv.normalize(hist_test, hist_test, alpha=0, beta=1, norm_type=cv.NORM_MINMAX)


    compare_method = cv.HISTCMP_CORREL

    base_base = cv.compareHist(hist_base, hist_base, compare_method)
    base_test = cv.compareHist(hist_base, hist_test, compare_method)



    


    if(round(base_test,5) >= 0.97):
        return False
    else:
        print('base_test Similarity = ', base_test)
        return True 

def getFrame(sec): 
    global old_frame,currentframe
    vidcap.set(cv2.CAP_PROP_POS_MSEC,sec*1000) 
    hasFrames,image = vidcap.read()
    if(not(os.path.exists('frames'))):
        os.makedirs('frames', exist_ok=False)
    if hasFrames: 
        if old_frame is not None:
            if(capture(old_frame,image)):
                name = 'frames/Frame(' + str(currentframe) + ').jpg'
                cv.imwrite(name, image)
                currentframe += 1
                old_frame=image
            else:
                name = 'frames/Frame(' + str(currentframe) + ').jpg'
                cv.imwrite(name, image)
                os.remove('frames/Frame(' + str(currentframe-1) + ').jpg') 
                currentframe += 1
                old_frame=image
                   
            
        else:
            name = 'frames/Frame(' + str(currentframe) + ').jpg'
            cv.imwrite(name, image)
            currentframe += 1
            old_frame=cv.imread('frames/Frame(0).jpg')
    return hasFrames 
def createppt():
        #Create presentation and setting layout as blank (6)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    #Find number of slides to create
    #First Method = Count number of images in screenshot files (change path depending on the user)
    # nbSlide = len(fnmatch.filter(os.listdir("frames"), '*.jpg'))

    for images in os.listdir("frames"):
        slide = prs.slides.add_slide(blank_slide_layout)
        #change background with an image of the slide …
        left = top = 0
        pic = slide.shapes.add_picture('frames/'+images, left-0.1*prs.slide_width, top, height = prs.slide_height)
    prs.save('final.pptx')
@csrf_exempt
def func(request): 
    global vidcap
    if request.method == 'POST':
        video = request.FILES['excel_file']
        vidcap=cv2.VideoCapture('E:/python/'+str(video)) 
        sec = 0 
        frameRate = 1.0 #it will capture image in each  second 
        success = getFrame(sec) 
        while success: 
            sec = sec + frameRate 
            sec = round(sec, 2) 
            success = getFrame(sec) 
        createppt()
        return render(request,"home.html")   

       